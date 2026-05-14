"""Build LearnAI CSCI 379 final presentation as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Palette ────────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
ACCENT    = RGBColor(0x6C, 0x63, 0xFF)   # violet
ACCENT2   = RGBColor(0x22, 0xD3, 0xEE)   # cyan
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0x94, 0xA3, 0xB8)
GREEN     = RGBColor(0x34, 0xD3, 0x99)
YELLOW    = RGBColor(0xFB, 0xBF, 0x24)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helpers ────────────────────────────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()  # no border
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def txbox(slide, text, left, top, width, height,
          size=28, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          italic=False, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def bullet_box(slide, items, left, top, width, height,
               size=20, color=WHITE, dot_color=ACCENT2, indent=Inches(0.35)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        # bullet dot
        dot = p.add_run()
        dot.text = "● "
        dot.font.size = Pt(size - 2)
        dot.font.color.rgb = dot_color
        # text
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(6)
    return tb


def checklist_box(slide, items, left, top, width, height, size=19, checked=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    mark = "✓ " if checked else "○ "
    mark_color = GREEN if checked else GREY
    first = True
    for label, done in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        tick = p.add_run()
        tick.text = "✓ " if done else "○ "
        tick.font.size = Pt(size)
        tick.font.color.rgb = GREEN if done else GREY
        run = p.add_run()
        run.text = label
        run.font.size = Pt(size)
        run.font.color.rgb = WHITE if done else GREY
        p.space_after = Pt(5)
    return tb


def accent_bar(slide, top=Inches(0.55), height=Pt(3)):
    """Thin coloured rule across the top."""
    add_rect(slide, Inches(0.5), top, Inches(12.33), int(height), ACCENT)


def slide_number(slide, n, total):
    txbox(slide, f"{n} / {total}",
          Inches(11.8), Inches(7.1), Inches(1.2), Inches(0.35),
          size=11, color=GREY, align=PP_ALIGN.RIGHT)


# ── Slide content ──────────────────────────────────────────────────────────────
TOTAL = 9


# 1. TITLE ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_bg(s)

# big violet rectangle left-half accent block
add_rect(s, Inches(0), Inches(0), Inches(5.5), SLIDE_H, RGBColor(0x1E, 0x1B, 0x4B))

# app name
txbox(s, "LearnAI", Inches(0.55), Inches(2.2), Inches(4.8), Inches(1.4),
      size=72, bold=True, color=ACCENT)
txbox(s, "AI-Powered Interactive Learning Platform",
      Inches(0.55), Inches(3.55), Inches(4.8), Inches(0.8),
      size=22, color=WHITE)
txbox(s, "CSCI 379  ·  Final Project  ·  May 2026",
      Inches(0.55), Inches(4.3), Inches(4.8), Inches(0.5),
      size=16, color=GREY)

# right-side tagline
txbox(s, "Enter any topic.\nGet a full story.\nLearn by doing.",
      Inches(6.3), Inches(2.8), Inches(6.5), Inches(2.0),
      size=34, color=WHITE)
txbox(s, "Stories  →  Chapters  →  Scenes  →  Quests",
      Inches(6.3), Inches(5.1), Inches(6.5), Inches(0.6),
      size=18, color=ACCENT2)

slide_number(s, 1, TOTAL)


# 2. WHAT IT DOES ──────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_bg(s)
accent_bar(s)

txbox(s, "What LearnAI Does", Inches(0.5), Inches(0.7), Inches(10), Inches(0.7),
      size=36, bold=True, color=WHITE)

# two columns
col1 = [
    'User enters any topic (e.g. "Data Structures")',
    "AI generates a structured story: 2 chapters, 2 scenes each",
    "Each scene has 3 quests — multiple choice, fill-in-the-blank, short answer",
    "Complete a scene to earn XP and unlock the next one",
]
col2 = [
    "Profile page tracks XP over time with a Chart.js line chart",
    "EN / ES language toggle (Gettext i18n)",
    "Google OAuth + magic-link email login",
    "Upload a .txt/.md reference doc to bias the AI prompt",
]

bullet_box(s, col1, Inches(0.5), Inches(1.6), Inches(5.9), Inches(4.5), size=19)
bullet_box(s, col2, Inches(6.7), Inches(1.6), Inches(5.9), Inches(4.5), size=19)

slide_number(s, 2, TOTAL)


# 3. LIVE DEMO ─────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_bg(s)

# big centered label
add_rect(s, Inches(2.5), Inches(2.0), Inches(8.3), Inches(2.8), RGBColor(0x1E, 0x1B, 0x4B))
txbox(s, "🎬  Live Demo", Inches(2.5), Inches(2.2), Inches(8.3), Inches(1.2),
      size=56, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

demo_steps = [
    "Register / log in with Google OAuth",
    "Create a new story — watch real-time generation progress",
    "Complete a scene (multiple choice → fill-blank → short answer)",
    "View XP chart on Profile",
    "Toggle EN ↔ ES language",
]
bullet_box(s, demo_steps, Inches(3.5), Inches(3.5), Inches(6.5), Inches(3.0),
           size=18, dot_color=YELLOW)

slide_number(s, 3, TOTAL)


# 4. MUST-HAVES CHECKLIST ──────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_bg(s)
accent_bar(s)

txbox(s, "Must-Haves  (19 pts)", Inches(0.5), Inches(0.7), Inches(10), Inches(0.7),
      size=36, bold=True, color=WHITE)

must_haves = [
    ("LiveView Auth — routes secured by Scope / on_mount", True),
    ("LiveView Real-time — PubSub story-generation progress bar", True),
    ("≥5 new functional components with attrs + styling", True),
    ("≥2 components replace JS behaviour (dark-mode toggle, lang toggle)", True),
    ("Transition animations — modals, navbar dropdowns", True),
    ("Dark Mode — functional toggle, >90% coverage", True),
    ("Mobile-first / breakpoints — no horizontal scroll at 320px", True),
    ("≥40% test coverage — all tests pass", True),
]

checklist_box(s, must_haves, Inches(0.5), Inches(1.55), Inches(12), Inches(5.5), size=20)

slide_number(s, 4, TOTAL)


# 5. OPTIONAL ITEMS ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_bg(s)
accent_bar(s)

txbox(s, "Chosen Optional Items  (7 × 3 pts = 21 pts)", Inches(0.5), Inches(0.7),
      Inches(12), Inches(0.7), size=34, bold=True, color=WHITE)

optional = [
    ("Google OAuth — Ueberauth + phx.gen.auth side-by-side", True),
    ("Associative Schemas — Story→Chapter→Scene→Quest (1:n chain);  User↔Scene via scene_completions (n:m)", True),
    ("Embedded Schemas — Quest.options uses embeds_many Quest.Option stored as JSONB", True),
    ("File Uploads — .txt / .md reference doc appended to AI prompt", True),
    ("Displaying Charts — Chart.js cumulative XP line chart on Profile", True),
    ("Mailer for Transactional Emails — welcome + magic-link via Swoosh (visible at /dev/mailbox)", True),
    ("Internationalization — ≥90% translated into Spanish via Gettext; EN/ES session toggle", True),
]

checklist_box(s, optional, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.2), size=19)

slide_number(s, 5, TOTAL)


# 6. ARCHITECTURE ──────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_bg(s)
accent_bar(s)

txbox(s, "Architecture Highlight", Inches(0.5), Inches(0.7), Inches(10), Inches(0.7),
      size=36, bold=True, color=WHITE)

# port/adapter diagram (text-art)
txbox(s, "AI Port / Adapter Pattern", Inches(0.5), Inches(1.6), Inches(5.5), Inches(0.5),
      size=20, bold=True, color=ACCENT2)

diagram = (
    "GeneratorPort  (behaviour)\n"
    "  generate_story/1\n"
    "  grade_answer/2\n"
    "        ↓\n"
    "  AI.ex  (façade — reads config)\n"
    "   ├── OpenAIAdapter  (prod)\n"
    "   └── StubAdapter    (dev/test)"
)
txbox(s, diagram, Inches(0.5), Inches(2.1), Inches(5.3), Inches(3.5),
      size=17, color=WHITE)

txbox(s, "PubSub Real-time Flow", Inches(6.7), Inches(1.6), Inches(5.5), Inches(0.5),
      size=20, bold=True, color=ACCENT2)

flow = (
    "Stories.create_story_async/2\n"
    "  → Task.Supervisor spawns task\n"
    "  → AI.generate_story/1\n"
    "  → insert chapters/scenes/quests\n"
    '  → PubSub.broadcast "story:{id}"\n'
    "        ↓\n"
    "StoryLive.New  handle_info/2\n"
    "  → progress bar update\n"
    "  → redirect on :story_ready"
)
txbox(s, flow, Inches(6.7), Inches(2.1), Inches(5.5), Inches(3.8),
      size=17, color=WHITE)

slide_number(s, 6, TOTAL)


# 7. LESSON LEARNED 1 (from class) ────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_bg(s)
accent_bar(s)

txbox(s, "Lesson Learned  #1", Inches(0.5), Inches(0.65), Inches(8), Inches(0.55),
      size=22, color=GREY)
txbox(s, "From class: LiveView is the server,\nnot the browser", Inches(0.5), Inches(1.2),
      Inches(12), Inches(1.4), size=38, bold=True, color=WHITE)

txbox(s, (
    "My first instinct for the quiz flow was to write JavaScript to manage state "
    "(current question, answer, phase). I built half of it before realising LiveView "
    "already owns the state on the server — I only needed to send events and re-render.\n\n"
    "Cutting the JS dropped ~80 lines and made the whole quiz flow testable with "
    "standard ExUnit + LiveView test helpers — no browser automation needed."
), Inches(0.5), Inches(2.8), Inches(12.3), Inches(3.5), size=21, color=WHITE)

txbox(s, "Takeaway: reach for a LiveView event before reaching for JavaScript.",
      Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
      size=19, bold=True, color=ACCENT2)

slide_number(s, 7, TOTAL)


# 8. LESSON LEARNED 2 (from project) ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_bg(s)
accent_bar(s)

txbox(s, "Lesson Learned  #2", Inches(0.5), Inches(0.65), Inches(8), Inches(0.55),
      size=22, color=GREY)
txbox(s, "From the project: isolate the thing\nyou can't control", Inches(0.5), Inches(1.2),
      Inches(12), Inches(1.4), size=38, bold=True, color=WHITE)

txbox(s, (
    "The OpenAI API is slow, expensive, and non-deterministic — impossible to test against "
    "directly. Hiding it behind a GeneratorPort behaviour meant I could write a StubAdapter "
    "that returns hardcoded Roman Empire data in milliseconds. Every context test, every "
    "LiveView test runs against the stub — zero API calls, zero flakiness.\n\n"
    "The same pattern would work for any external dependency: payment provider, "
    "third-party API, even the database. Define a behaviour first; the real adapter is "
    "just one implementation of it."
), Inches(0.5), Inches(2.8), Inches(12.3), Inches(3.5), size=21, color=WHITE)

txbox(s, "Takeaway: one behaviour + two adapters turns an external dependency into a seam you control.",
      Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
      size=19, bold=True, color=ACCENT2)

slide_number(s, 8, TOTAL)


# 9. CLOSING ───────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_bg(s)

add_rect(s, Inches(0), Inches(0), Inches(5.5), SLIDE_H, RGBColor(0x1E, 0x1B, 0x4B))

txbox(s, "Thank you", Inches(0.55), Inches(2.6), Inches(4.8), Inches(1.0),
      size=56, bold=True, color=ACCENT)
txbox(s, "Questions?", Inches(0.55), Inches(3.65), Inches(4.8), Inches(0.7),
      size=28, color=WHITE)

summary = [
    "19 / 19  must-have points",
    "21 / 21  optional points",
    "Full test suite passing",
    "Real AI  ·  Real emails  ·  Real OAuth",
]
bullet_box(s, summary, Inches(6.3), Inches(2.4), Inches(6.5), Inches(3.0),
           size=24, dot_color=GREEN)

slide_number(s, 9, TOTAL)


# ── Save ───────────────────────────────────────────────────────────────────────
out = "/home/hqp001/Projects/csci379_final/slides/LearnAI_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}")
